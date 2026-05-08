"""
Generate PowerPoint slide deck from agent execution audit trail.

This script reads the agent_state.db SQLite checkpoint and creates
a professional presentation summarizing the agent's execution,
findings, and decisions.

Usage: python generate_slides.py [--thread-id RUN_ID] [--output slides.pptx]
"""

import json
import sqlite3
from datetime import datetime
from pathlib import Path

# For now, just a skeleton. To implement:
# 1. Install: pip install python-pptx
# 2. OR use MCP: office-powerpoint-mcp-server (https://github.com/GongRzhe/Office-PowerPoint-MCP-Server)


def get_audit_trail(db_path: str = "agent_state.db", thread_id: str | None = None) -> dict:
    """
    Retrieve agent execution history from SQLite checkpoint.

    Args:
        db_path: Path to agent_state.db
        thread_id: Optional thread_id to filter specific runs

    Returns:
        Dictionary with execution metadata and step results
    """
    conn = sqlite3.connect(db_path)
    cursor = conn.cursor()

    # Query the checkpoint data
    # LangGraph stores state in 'checkpoint' table with JSON
    try:
        cursor.execute("""
            SELECT thread_id, checkpoint FROM checkpoint
            ORDER BY ts_created DESC
            LIMIT 1
        """)
        row = cursor.fetchone()

        if row:
            thread_id_db, checkpoint_json = row
            checkpoint_data = json.loads(checkpoint_json)
            return {
                "thread_id": thread_id_db,
                "timestamp": datetime.now().isoformat(),
                "execution": checkpoint_data,
            }
    except Exception as e:
        print(f"Error reading checkpoint: {e}")

    conn.close()
    return {}


def generate_slides_with_pptx(audit_data: dict, output_path: str = "output/agent_report.pptx") -> None:
    """
    Generate PowerPoint using python-pptx library.

    Requires: pip install python-pptx

    This is a template — fill in slides based on audit_data.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN

        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        # Slide 1: Title
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = "Multi-Step Agent Execution Report"
        subtitle.text = f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"

        # Slide 2: Execution Summary
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        title = slide.shapes.title
        body_shape = slide.placeholders[1]

        title.text = "Execution Summary"
        tf = body_shape.text_frame
        tf.text = "Agent completed task successfully"

        execution = audit_data.get("execution", {})
        if isinstance(execution, dict):
            for key, value in list(execution.items())[:3]:
                p = tf.add_paragraph()
                p.text = f"{key}: {str(value)[:50]}"
                p.level = 1

        # Slide 3: Steps Executed
        slide = prs.slides.add_slide(bullet_slide_layout)
        title = slide.shapes.title
        body_shape = slide.placeholders[1]

        title.text = "Steps Executed"
        tf = body_shape.text_frame
        tf.text = "1. Data Collection (SQL query)"
        for step_num in range(2, 5):
            p = tf.add_paragraph()
            p.text = f"{step_num}. Analysis & Report"
            p.level = 0

        # Save
        Path(output_path).parent.mkdir(parents=True, exist_ok=True)
        prs.save(output_path)
        print(f"✓ Slide deck created: {output_path}")

    except ImportError:
        print(
            "python-pptx not installed. Install with: pip install python-pptx\n"
            "Or use the Office PowerPoint MCP Server instead:\n"
            "  uvx office-powerpoint-mcp-server"
        )


def generate_slides_with_mcp(audit_data: dict) -> None:
    """
    Generate PowerPoint using Office-PowerPoint-MCP-Server MCP.

    This is ideal because it integrates with Claude Code.
    Requires: pip install office-powerpoint-mcp-server
    Or: uvx office-powerpoint-mcp-server

    In Claude Code, ask: "Generate slides from this audit trail..."
    """
    print(
        "To use the MCP approach:\n"
        "1. Install: uvx office-powerpoint-mcp-server\n"
        "2. Configure in Claude Code settings.json\n"
        "3. Ask Claude Code: 'Generate slides from this audit trail...'\n"
        "\nMCP Server: https://github.com/GongRzhe/Office-PowerPoint-MCP-Server"
    )


def main():
    """Generate slide deck from agent execution."""
    import argparse

    parser = argparse.ArgumentParser(description="Generate PowerPoint from agent audit trail")
    parser.add_argument("--db", default="agent_state.db", help="Path to agent_state.db")
    parser.add_argument("--thread-id", help="Specific thread_id to use")
    parser.add_argument("--output", default="output/agent_report.pptx", help="Output .pptx path")
    parser.add_argument("--method", choices=["pptx", "mcp"], default="pptx", help="Generation method")

    args = parser.parse_args()

    print("📊 Generating slide deck from agent execution...")
    audit_data = get_audit_trail(args.db, args.thread_id)

    if not audit_data:
        print("⚠ No execution data found. Run 'python main.py' first.")
        return

    print(f"✓ Loaded audit trail: {audit_data.get('timestamp')}")

    if args.method == "pptx":
        generate_slides_with_pptx(audit_data, args.output)
    else:
        generate_slides_with_mcp(audit_data)


if __name__ == "__main__":
    main()
